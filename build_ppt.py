"""Generate the Titus VP presentation deck (HVAC AI Takeoff Tool)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

NAVY = RGBColor(0x0E, 0x2A, 0x47)
BLUE = RGBColor(0x1F, 0x6F, 0xB2)
GREY = RGBColor(0x44, 0x4444 // 256, 0x44)  # placeholder; fixed below
GREY = RGBColor(0x55, 0x55, 0x55)
GREEN = RGBColor(0x1E, 0x8E, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def title_slide(title, subtitle, tag):
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
    bg.shadow.inherit = False
    tb = s.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.5), Inches(2.5)).text_frame
    tb.word_wrap = True
    p = tb.paragraphs[0]; r = p.add_run(); r.text = title
    r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = WHITE
    p2 = tb.add_paragraph(); r2 = p2.add_run(); r2.text = subtitle
    r2.font.size = Pt(20); r2.font.color.rgb = RGBColor(0xBF, 0xD7, 0xEA)
    p3 = tb.add_paragraph(); r3 = p3.add_run(); r3.text = tag
    r3.font.size = Pt(14); r3.font.color.rgb = RGBColor(0x9F, 0xB3, 0xC8)
    return s


def content_slide(title, bullets, notes, kicker=None):
    s = prs.slides.add_slide(BLANK)
    bar = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.15))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    bar.shadow.inherit = False
    t = s.shapes.add_textbox(Inches(0.6), Inches(0.22), Inches(12.1), Inches(0.8)).text_frame
    t.word_wrap = True
    rp = t.paragraphs[0]; rr = rp.add_run(); rr.text = title
    rr.font.size = Pt(28); rr.font.bold = True; rr.font.color.rgb = WHITE
    if kicker:
        k = s.shapes.add_textbox(Inches(0.6), Inches(1.25), Inches(12.1), Inches(0.5)).text_frame
        k.word_wrap = True
        kr = k.paragraphs[0].add_run(); kr.text = kicker
        kr.font.size = Pt(15); kr.font.italic = True; kr.font.color.rgb = BLUE
    body = s.shapes.add_textbox(Inches(0.7), Inches(1.75 if kicker else 1.45),
                                Inches(12.0), Inches(5.4)).text_frame
    body.word_wrap = True
    for i, (text, lvl) in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.level = lvl
        run = p.add_run(); run.text = ("• " if lvl == 0 else "– ") + text
        run.font.size = Pt(19 - 2 * lvl)
        run.font.color.rgb = NAVY if lvl == 0 else GREY
        if lvl == 0:
            run.font.bold = True
        p.space_after = Pt(7)
    s.notes_slide.notes_text_frame.text = notes
    return s


# ---- slides ----
title_slide(
    "HVAC AI Takeoff Tool",
    "Automated equipment takeoff from construction drawings — detect, spec, and quantify HVAC air-distribution & equipment",
    "Triune Solutions  ·  Presented to the VP of Titus  ·  June 2026")

content_slide("The problem we set out to solve",
  [("HVAC estimators do takeoffs by hand — slow, costly, error-prone", 0),
   ("Count every diffuser, grille, register, damper, fan & unit across dozens of E-size sheets", 1),
   ("Look each tag up in the equipment schedule for its specs (CFM, neck size, model, brand)", 1),
   ("Record it, mark it off, build the Bill of Materials", 1),
   ("A single project = hours to days of manual counting; mistakes cost bids and margin", 0)],
  "Frame it from the estimator's pain. A commercial HVAC project has dozens of 36x24 sheets, each packed "
  "with hundreds of symbols. The estimator manually counts each device, cross-references the schedule for "
  "specs, and types it into Excel. It's tedious, slow, and easy to miss items or mis-spec. That directly "
  "affects bid accuracy and profit. This is the job we're automating.")

content_slide("What we built",
  [("An AI tool that reads a drawing-set PDF and produces a quantified takeoff", 0),
   ("Detects equipment symbols on the plans with computer vision (YOLOv8)", 1),
   ("Matches each to its tag and the equipment schedule (specs)", 1),
   ("Reconciles detected counts against the schedule; flags discrepancies", 1),
   ("Outputs an Excel takeoff + an annotated & Bluebeam-ready PDF — in the team's format", 1),
   ("Learns continuously from estimator corrections", 0)],
  "Plain summary: feed it the PDF, it gives back a takeoff. Five steps — ingest, detect, tag/schedule, "
  "reconcile, output — plus a learning loop. It's a working web app the estimators use: upload, view "
  "detections on the plan, correct mistakes, download the Excel. Stack: FastAPI backend + Next.js frontend, "
  "YOLOv8 for vision, PyMuPDF + OCR for the document, openpyxl for Excel.")

content_slide("How it works — the pipeline",
  [("Part 1 · Intake & Page Understanding — pick the real plan pages, read legend/schedule/title block", 0),
   ("Part 2 · Detection — YOLO finds the equipment symbols on those plans", 0),
   ("Part 3 · Extraction & Reconciliation — parse schedule, assign tags, match & reconcile", 0),
   ("Part 4 · Delivery & Learning — Excel/Bluebeam output + the self-improving loop", 0),
   ("Each part owns one job and one accuracy number — so we can fix and measure them independently", 0)],
  "We split the system into 4 parts by the accuracy each one owns. This matters: it lets us measure and "
  "improve each independently instead of guessing. Part 1 decides WHERE to look; Part 2 decides WHAT'S there; "
  "Part 3 reads the specs and matches; Part 4 delivers and learns. Key point for the VP: only Part 2 is a "
  "trained AI model — the rest is deterministic rules, which makes them reliable and auditable.")

content_slide("Part 1 — Reading the document (page selection)",
  [("Picks the real floor-plan pages across ANY engineer's drawing style", 0),
   ("Reads each sheet's number/title (text + OCR), classifies plan vs schedule/legend/detail", 1),
   ("Routes pages: plans → detection, schedules → spec parsing, legend → symbol dictionary", 1),
   ("Result: zero real plans dropped (recall 1.00 on measured styles)", 0),
   ("Recovered plans the old logic skipped (one set: 4 → 10 pages) and removed phantom pages", 1)],
  "Garbage-in prevention. If you scan the wrong pages you miss equipment or invent phantoms. We replaced a "
  "brittle filter with a fused decision that votes across multiple signals, so it's robust to different "
  "firms' numbering/title styles. Concrete wins: one project went from 4 to 10 detected plan pages (it was "
  "silently dropping 6 real plans); another stopped counting a schedule sheet as a plan (a phantom source). "
  "Measured recall = 1.00 — we don't drop real plans.")

content_slide("Part 2 — Detecting the equipment (the AI's eyes)",
  [("YOLOv8 vision model finds diffusers, grilles, registers, dampers, fans, units on the plans", 0),
   ("Measured on held-out drawings the model never trained on:", 0),
   ("Object-level recall 84%, precision 94% — it reliably finds the air-distribution products", 1),
   ("Key insight: the model FINDS the device; the SCHEDULE supplies the exact type & specs", 0),
   ("So vision + schedule together = a complete, specced count — without guessing from the symbol", 1)],
  "This is the headline and it's Titus-relevant. The raw benchmark first read low (~31%), but we proved "
  "that was mostly a LABELING mismatch, not blindness — the model actually detects ~84% of the air devices "
  "on held-out plans, at 94% precision. Diffusers/grilles look nearly identical, so the model can't tell "
  "the subtype visually — but it doesn't need to: the tag + schedule tell us the exact product, neck size, "
  "and model. That division of labor is the core design and it's exactly the air-distribution category "
  "Titus makes.")

content_slide("Part 3 — Schedules, tags & reconciliation",
  [("Parses the equipment schedule: each tag → CFM, neck size, model, brand, mounting", 0),
   ("Assigns each detected symbol its tag (direct map, callout text, bubble OCR)", 0),
   ("Matches detections to schedule rows; outputs specced line items", 0),
   ("Reconciles counts: 'schedule lists 7 RTUs, we found 5' → flags the gap honestly", 0),
   ("Product names come straight from the schedule (e.g. LAY-IN, PERFORATED FACE)", 0)],
  "This is where it becomes a real takeoff, not just boxes. We read the schedule into per-tag specs, attach "
  "each detection to its tag, and produce line items with brand/model/neck/CFM. Reconciliation is the trust "
  "layer — it compares what we found vs what the schedule says should exist and flags differences instead of "
  "hiding them. And the product description in the output comes from the schedule, matching how the team "
  "already writes takeoffs.")

content_slide("Part 4 — Output & the self-improving loop",
  [("Outputs: Excel takeoff (team's exact format) + annotated PDF + Bluebeam count stamps", 0),
   ("Estimator corrects mistakes in the UI / Bluebeam — every fix becomes training data", 0),
   ("Corrections → retrain the model → benchmark gate → deploy if it's better", 0),
   ("A flywheel: more use → more corrections → smarter model → less correction needed", 0),
   ("Safety rule: never ship a model that doesn't beat the current one on held-out drawings", 0)],
  "The output is drop-in for the team's workflow (byte-identical Excel format; Bluebeam stamps they work "
  "with). The strategic part is the flywheel: the tool improves from normal work, not a separate labeling "
  "project. The benchmark 'gate' is the discipline that keeps quality up — a new model only ships if it "
  "objectively beats the old one on drawings it never saw. We caught and rejected a bad model this way "
  "(it would have regressed accuracy).")

content_slide("Accuracy & trust — how we keep it honest",
  [("We measure against the team's own completed takeoffs and Bluebeam-marked ground truth", 0),
   ("Page selection: recall 1.00 (no dropped plans)", 1),
   ("Detection: 84% object recall / 94% precision on held-out plans", 1),
   ("The 'gate' blocks any model that doesn't improve on held-out data", 0),
   ("Reconciliation surfaces misses instead of hiding them — estimator stays in control", 0)],
  "Trust is everything for adoption. Two principles: (1) measure on held-out drawings the model never saw — "
  "real generalization, not memorization; (2) the tool is an ASSIST — it shows its work, flags what it's "
  "unsure of, and the estimator verifies. We don't claim 100%; we claim honest, measurable, and improving. "
  "That's how an estimating team actually trusts and uses it.")

content_slide("The data advantage",
  [("We mined Triune's own completed takeoffs into training data", 0),
   ("279 past projects → ~46,000 labeled training tiles", 1),
   ("Equipment the model used to miss is now well-represented:", 0),
   ("Rooftop units 6 → 435 examples · Fire-smoke dampers 6 → 823 · Air devices 1,036 → 54,249", 1),
   ("The model learns from real Triune work — a compounding, proprietary asset", 0)],
  "This is a durable competitive moat. Years of completed takeoffs are a goldmine of labeled examples. We "
  "built the pipeline to mine them automatically. The classes the model was blind to (because it had ~6 "
  "examples) now have hundreds-to-thousands. Every project the team does makes the data — and the model — "
  "stronger. No competitor without this data history can replicate it.")

content_slide("Where it stands today",
  [("Live web app: upload → detect → correct → download takeoff", 0),
   ("Working now (production model): page selection, detection, schedule matching, correct product naming, honest reconciliation", 0),
   ("Built & validated end-to-end: the retrain loop (data, recipe, benchmark gate)", 0),
   ("One step from a major detection upgrade: a ~30-min GPU training run on the new data", 0),
   ("Roadmap-ready: the architecture is modular and measured", 0)],
  "Be honest and confident. What's live today already helps: right pages, ~84% detection, correct specs "
  "from the schedule, fewer false positives, honest flagging. The next jump — detecting the rarer equipment "
  "(RTUs, fire-smoke dampers) — is fully built and one GPU training run from deploying; the data and the "
  "safety gate are ready. We're not pitching vapor; we're pitching a working tool with a clear, funded next step.")

content_slide("Why this matters to Titus",
  [("The tool's strongest capability is finding & speccing air-distribution products — Titus's core", 0),
   ("It already extracts neck size, CFM, model, mounting per device", 0),
   ("Natural bridge to Titus Selection Navigator → part numbers, selections, pricing", 0),
   ("Faster, more accurate takeoffs → more Titus product correctly specified & ordered", 0),
   ("A data partnership could tune the model to Titus product lines specifically", 0)],
  "This is the slide the VP cares about most. Our 84% sweet spot IS air distribution — diffusers, grilles, "
  "registers, louvers — Titus's bread and butter. We already pull neck size, CFM, and model from the "
  "schedule, which is exactly what Titus Selection Navigator needs as input. Imagine: drawing in → "
  "Titus-product BOM out, with selections and pricing. That accelerates quoting and pulls more Titus product "
  "through. There's a real partnership/integration story here, not just a software demo.")

content_slide("Future prospects",
  [("Close the retrain loop → higher recall across all equipment classes", 0),
   ("Titus integration: detected air devices → Selection Navigator → auto Titus BOM + pricing", 0),
   ("Multi-document: spec books, addenda, Division 23 — full-project understanding", 0),
   ("Expand beyond HVAC: plumbing & electrical takeoffs on the same engine", 0),
   ("SaaS product for HVAC contractors & suppliers (a Rebar-class competitor)", 0)],
  "Paint the trajectory. Near term: finish the detection upgrade, ship it. Mid term: the Titus integration "
  "(the differentiator) and multi-document understanding. Long term: extend the same vision+schedule engine "
  "to plumbing and electrical, and productize as SaaS for the broader market. Each step compounds on the "
  "data flywheel. The Titus angle can be the wedge that funds and shapes the roadmap.")

content_slide("The vision",
  [("From: estimator manually counts and specs every device for hours", 0),
   ("To: drawing in → instant, specced, manufacturer-ready takeoff", 0),
   ("Faster quoting · fewer errors · more product correctly specified", 0),
   ("An assistant that gets smarter with every project the team completes", 0)],
  "Land the plane on the big picture. We're turning a slow, manual, error-prone job into an instant, "
  "accurate, self-improving one — and the output is specced to the manufacturer level. For Titus that means "
  "more of the right product specified, faster. For Triune it's a proprietary, compounding asset. Close on "
  "the partnership opportunity.")

content_slide("Discussion & next steps",
  [("Demo the live tool on a real Titus-heavy project", 0),
   ("Explore a Titus Selection Navigator integration (product mapping & pricing)", 0),
   ("Data collaboration to tune the model to Titus product lines", 0),
   ("Define a pilot: scope, success metrics, timeline", 0)],
  "End with concrete asks, scaled to how the meeting goes. Minimum: offer a live demo. Bigger: propose the "
  "Selection Navigator integration and a pilot. Have a clear, low-friction next step ready so the VP can say "
  "yes to something.")

prs.save("HVAC_AI_Takeoff_Titus.pptx")
print("saved HVAC_AI_Takeoff_Titus.pptx |", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
